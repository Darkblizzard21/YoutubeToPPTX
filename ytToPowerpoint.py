import argparse
from datetime import datetime
import os
from os.path import abspath
from pptx import Presentation
from pptx.presentation import Presentation as PresentationType
from pptx.util import Cm 
from pptx.enum.shapes import PP_MEDIA_TYPE
from lxml import etree
import yt_dlp

global verbose 
global picturesOnly
global cacheFolder

def verbosePrint(msg):
    if verbose:
        print(msg)

def downloadVideo(code: str, videoFormat: str):
    fileName = str(os.path.join(rootFolder,"cache", code+"-"+videoFormat+".mp4"))
    
    if os.path.exists(fileName):
        print("Video of code: " + code + " is already present, skipped download.")
        return fileName

    yt_opts = {
        'outtmpl': fileName,
        'verbose': verbose,
        "format": videoFormat
    }

    with yt_dlp.YoutubeDL(yt_opts) as ydl:
        ydl.download("https://www.youtube.com/watch?v="+code) 
    return fileName


def makeVideoClip(videoPath, startTime, endTime):
    fileName = str(startTime).replace(":",".")  + "-" + str(endTime).replace(":",".") + ".mp4"
    filePath = str(os.path.join(cacheFolder, fileName))
    
    if os.path.exists(filePath):
        return filePath
    
    args = ["ffmpeg", "-y"]
    if not verbose:
        args += ["-loglevel", "error"]
    args += ["-ss", startTime]
    if endTime:
        args +=["-to", endTime]
    args += ["-i", "\""+ videoPath +"\""]
    args.append("\""+filePath+"\"")

    command = " ".join(args)
    print("Invoking: " + command)
    returnCode = os.system(command)
            
    if returnCode != 0:
        print("Process exited with code: " + str(returnCode))
        return None
    
    print("Successfully created clip: " + fileName)
    return filePath

def makeThumbnail(videoPath, timeStamp):
    fileName = str(timeStamp).replace(":",".") + ".jpg"
    filePath = str(os.path.join(cacheFolder, fileName))
    
    if os.path.exists(filePath):
        return filePath
    
    args = ["ffmpeg -y "]
    if not verbose:
        args += ["-loglevel error "]
    if timeStamp:
        args += ["-ss ", timeStamp, " -i \"", videoPath, "\" -frames:v 1 \"", filePath, "\""]
    else:
        # time stamp is none just give the last frame of the video
        args += ["-sseof -0.3 -i \"", videoPath, "\" -qscale:v 2 -update 1 \"", filePath, "\""]

    command = "".join(args)
    print("Invoking: " + command)
    returnCode = os.system(command)
            
    if returnCode != 0:
        print("Process exited with code: " + str(returnCode))
        return None
    
    print("Successfully created thumbnail: " + fileName)
    return filePath

def addPictureSlideSlide(presentation: PresentationType, picture): 
    blank_slide = presentation.slides.add_slide(presentation.slide_layouts[6]) 
    blank_slide.shapes.add_picture(picture, 
                                   left = Cm(0), top = Cm(0), 
                                   width = presentation.slide_width, height = presentation.slide_height)
    return blank_slide

def addVideoSlide(presentation: PresentationType, video, thumbnail): 
    if picturesOnly:
        return addPictureSlideSlide(presentation,thumbnail)

    blank_slide = presentation.slides.add_slide(presentation.slide_layouts[6]) 
    movie_shape = blank_slide.shapes.add_movie(video, 
                                               left = Cm(0), top = Cm(0), 
                                               width = presentation.slide_width, height = presentation.slide_height, 
                                               poster_frame_image = thumbnail, mime_type = "video/mp4")
    assert movie_shape.media_type == PP_MEDIA_TYPE.MOVIE
    tree = movie_shape._element.getparent().getparent().getnext().getnext()
    timing = [el for el in tree.iterdescendants() if etree.QName(el).localname == 'cond'][0]
    timing.set('delay', '0')
    return blank_slide

class Instruction:
  def __init__(self, timeStamp, instruction, args = None):
    self.timeStamp = timeStamp
    self.instruction = instruction
    self.args = args

if __name__ == "__main__":
    start_time = datetime.now()                 
    rootFolder = os.path.dirname(__file__)

    parser = argparse.ArgumentParser(
                        prog='Powerpoint Form Video',
                        description='Build a powerpoint from a video',
                        epilog='Written by Pirmin Pfeifer',
                        formatter_class=argparse.RawTextHelpFormatter)

    parser.add_argument('-i', "--input",
                        default=os.path.join(rootFolder, 'presentation.md'),
                        required=False,
                        help="Presentation Manifest used to generate presentation")
    parser.add_argument('-t','--template',
                        default=os.path.join(rootFolder, 'Template.pptx'),
                        required=False,
                        help="Template powerpoint file, the generated presentation will be appended to the file") 
    parser.add_argument('-f','--format',
                        default=136,
                        required=False,
                        help="Specify video format to use (default is 136 (720P)"+
                        "\nOther common formats are: 401 (2160P), 400 (1440P), 136 (720P), 397 (480P), 396 (360P)"+
                        "\nNote selecting a format that is not provided will result in the script crashing non gracefully") 

    parser.add_argument('-v','--verbose',
                        action='store_true') 
    parser.add_argument('-s','--short',
                        action='store_true',
                        help="Will produce a short pptx just containing the last frame of each specified clip") 
    parser.add_argument('-r','--reverse',
                        action='store_true',
                        help="Development option: Reverses slide order for faster iteration") 
    parser.add_argument('--no-output',
                        action='store_true',
                        help="Development option: Skips saving the final pptx to the hard drive.") 
    args = parser.parse_args()

    verbose = args.verbose
    picturesOnly = args.short
    videoFormat = str(args.format)
    # def options
    reverse = args.reverse
    no_output = args.no_output

    manifest = abspath(args.input)
    if not os.path.exists(manifest):
        print("Presentation Manifest: " + manifest + " dose not exist!")
        exit()

    template = abspath(args.template)
    prs = Presentation()
    if os.path.exists(template):
        prs = Presentation(template)
    else:
        print("Template was not found ... continuing without")

    videoFile = ""
    instructions = [ Instruction("00:00:00","/skip")]
    with open(manifest) as f:
        code = f.readline().rstrip()
        cacheFolder = os.path.join(rootFolder, "cache", code, videoFormat)
        if not os.path.exists(cacheFolder):
            os.makedirs(cacheFolder) 
            
        videoFile = downloadVideo(code, videoFormat)

        while line := f.readline():
            line = line.rstrip()
            if line[0] == '#' :
                continue
            lineParts = line.split(' ')
            args = lineParts[2:]
            if len(args) == 0:
                args = None
            instructions.append(Instruction(lineParts[0],lineParts[1], args))
            
    instructions.append(Instruction(None,"/skip"))
    verbosePrint("Starting to create " + str(len(instructions)) + " slides.")

    indices = range(0, len(instructions)-1)
    if(reverse):
        indices = reversed(indices)
    for index in indices:
        instruction     = instructions[index]
        nextInstruction = instructions[index+1]

        code = instruction.instruction
        if code == "/copy" or code == "/note":
            clip = makeVideoClip(videoFile, instruction.timeStamp, nextInstruction.timeStamp)
            if not clip:
                print("ERROR while creating clip!")
                exit()
            thumb = makeThumbnail(videoFile, nextInstruction.timeStamp if picturesOnly else instruction.timeStamp)
            
            if not clip:
                print("ERROR while creating thumbnail!")
                exit()
            verbosePrint("copy " + instruction.timeStamp + "-" + str(nextInstruction.timeStamp))
            
            if no_output:
                continue
            slide = addVideoSlide(prs, clip, thumb)
            if code == "/note":
                verbosePrint(" also add note")
                notes_slide = slide.notes_slide
                text_frame = notes_slide.notes_text_frame
                text_frame.text = " ".join(instruction.args)
        elif code == "/skip":
            verbosePrint("skip")
        else:
            print("unknown instruction \"" + code + "\" at " + str(index) + " is treated as /skip")


    if no_output:
        print("Skipped output due to --no-output flag. No pptx file was written to drive.")
        print("Script finished in %s" % (datetime.now() - start_time))
        exit()

    extension = ".pptx" if not picturesOnly else ".short.pptx"
    output = os.path.join(rootFolder, os.path.splitext(os.path.basename(manifest))[0]+ extension)
    prs.save(output)
    print("Successfully generated: %s in %s" % (output, datetime.now() - start_time))

    


